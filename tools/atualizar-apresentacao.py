# -*- coding: utf-8 -*-
"""
atualizar-apresentacao.py — atualiza cirurgicamente a Apresentacao-Metodo-Norte.pptx

NÃO recria os slides (o visual foi feito à mão). Abre o .pptx existente e substitui
apenas os textos que mudam entre versões, preservando fontes, cores e layout.

Uso:  python tools/atualizar-apresentacao.py   (na raiz do kit)

Cada substituição é (texto_antigo -> texto_novo). Se o texto antigo não for achado,
o script AVISA (para não passar despercebido numa próxima versão da apresentação).
Edite as substituições abaixo a cada release.
"""
import sys
from pptx import Presentation

# console do Windows costuma ser cp1252 — forçar UTF-8 na saída p/ não quebrar em acentos
try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

ARQ = "Apresentacao-Metodo-Norte.pptx"

# (antigo, novo) — casam contra o texto COMPLETO do parágrafo, substituindo run a run.
# Usa-se travessão "–" ou "—"? A apresentação usa o separador "·" no rodapé — copiar EXATO.
SUBSTITUICOES = [
    # Rodapé (slide 11) — versão v1.3.0 + mudança de licença: proprietária -> MIT open-source
    (
        "Método Norte v1.1 · © 2026 Kasuo · Todos os direitos reservados",
        "Método Norte v1.3.0 · © 2026 Kasuo · Open-source (Licença MIT)",
    ),
    # Slide 8 — cartão de licença: refletir MIT + créditos no lugar de "licença de uso"
    (
        "Licença de uso claro, formulário de feedback e processo de melhoria contínua",
        "Open-source sob Licença MIT, com créditos a ideias externas adotadas e processo de melhoria contínua",
    ),
]


def substituir_paragrafo(paragraph, antigo, novo):
    """Se o texto completo do parágrafo == antigo, reescreve preservando o 1º run."""
    full = "".join(r.text for r in paragraph.runs)
    if full != antigo:
        return False
    # coloca todo o texto novo no primeiro run e esvazia os demais (mantém formatação do 1º)
    if paragraph.runs:
        paragraph.runs[0].text = novo
        for r in paragraph.runs[1:]:
            r.text = ""
    return True


def main():
    prs = Presentation(ARQ)
    pendentes = {a for a, _ in SUBSTITUICOES}
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for antigo, novo in SUBSTITUICOES:
                    if substituir_paragrafo(para, antigo, novo):
                        pendentes.discard(antigo)
                        print(f'  [OK] substituido: "{antigo[:50]}..."')

    if pendentes:
        print("\n[AVISO] NAO encontrados (texto pode ter mudado numa versao nova do .pptx):")
        for p in pendentes:
            print(f'    - "{p[:60]}..."')
        print("   Revise as SUBSTITUICOES no script antes de confiar no resultado.")
        sys.exit(1)

    prs.save(ARQ)
    print(f"\n[OK] {ARQ} atualizado ({len(SUBSTITUICOES)} substituicoes).")


if __name__ == "__main__":
    main()
