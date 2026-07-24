# -*- coding: utf-8 -*-
"""
adicionar-slide-relatorios.py — insere UM slide novo ("Relatórios que a IA devolve")
na Apresentacao-Metodo-Norte.pptx, no mesmo padrão visual dos slides de cartões.

Idempotente: se o slide já existir (título encontrado), não duplica.
Posição: logo após o slide 5 ("Os 8 pilares do método").

Uso:  python tools/adicionar-slide-relatorios.py   (na raiz do kit, .pptx fechado)
"""
import copy
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

ARQ = "Apresentacao-Metodo-Norte.pptx"
TITULO = "Os relatórios que a IA devolve"

# paleta extraída da própria apresentação
BG      = RGBColor(0x1E, 0x27, 0x61)  # fundo do slide
CARD    = RGBColor(0x27, 0x33, 0x6F)  # preenchimento do cartão
GOLD    = RGBColor(0xF0, 0xB4, 0x29)  # título do cartão
LIGHT   = RGBColor(0xCA, 0xDC, 0xFC)  # descrição
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# 4 cartões: 2 arquivos-relatório vivos + 2 check-ins que a IA devolve
CARDS = [
    ("ESTADO-ATUAL", "Onde o projeto está agora e o que falta — as pendências abertas, segurança em destaque. O relatório de 1 página que a IA lê e reporta ao abrir a sessão."),
    ("Progresso mensal", "O histórico do que foi feito, quando e por quê — fatiado por mês para não pesar no contexto. Cada tarefa concluída vira uma entrada curta."),
    ("Plano antes de agir", "Antes de tocar no código, a IA devolve um relatório: sistema, arquivos lidos, área impactada, risco e ação planejada. Você aprova antes."),
    ("Verificação provada", "Ao concluir, a IA relata o resultado observado (não deduzido) e, ao corrigir bugs, os gêmeos varridos (TWINS). Confiança com evidência."),
]

# layout de grade (mesma métrica do slide de pilares: 4 colunas, cartão 2.85 x 3.9)
LEFTS = [0.70, 3.75, 6.80, 9.85]
CARD_TOP = 1.90
CARD_W, CARD_H = 2.85, 4.35


def add_textbox(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def main():
    prs = Presentation(ARQ)

    # idempotência: já existe?
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and TITULO in sh.text_frame.text:
                print(f"[SKIP] slide '{TITULO}' já existe — nada a fazer.")
                return

    layout = prs.slides[4].slide_layout  # mesmo layout DEFAULT dos cartões
    slide = prs.slides.add_slide(layout)

    # fundo escuro igual aos demais
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # título
    t = add_textbox(slide, 0.70, 0.50, 12.0, 0.9, TITULO, 34, WHITE, bold=True)
    t.text_frame.paragraphs[0].runs[0].font.name = "Cambria"

    # subtítulo
    add_textbox(slide, 0.70, 1.25, 12.0, 0.5,
                "A documentação não é estática: ela relata. Rastreabilidade do início ao fim de cada tarefa.",
                14, LIGHT)

    # cartões
    for (titulo, desc), left in zip(CARDS, LEFTS):
        card = slide.shapes.add_shape(1, Inches(left), Inches(CARD_TOP), Inches(CARD_W), Inches(CARD_H))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD
        card.line.fill.background()
        card.shadow.inherit = False
        # título do cartão (dourado)
        add_textbox(slide, left + 0.15, CARD_TOP + 0.25, CARD_W - 0.3, 0.6, titulo, 15, GOLD, bold=True)
        # descrição (azul-claro)
        add_textbox(slide, left + 0.15, CARD_TOP + 0.95, CARD_W - 0.3, CARD_H - 1.1, desc, 11.5, LIGHT)

    # mover o slide novo para logo após o slide 5 (índice 5)
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    novo = slides[-1]
    xml_slides.remove(novo)
    xml_slides.insert(5, novo)  # posição 6 (após "Os 8 pilares")

    prs.save(ARQ)
    print(f"[OK] slide '{TITULO}' inserido como slide 6.")


if __name__ == "__main__":
    main()
